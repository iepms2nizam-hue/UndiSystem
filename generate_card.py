from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
import os

# Fungsi samarkan IC
def mask_ic(ic):
    return ic[:2] + "****" + ic[6:8] + "***" + ic[-1:]

# Data client
client_data = {
    "nama": "Razimi bin Komal",
    "ic": "702101126565",
    "whatsapp": "01123456789",
    "email": "razimi@example.com",
    "id": "6"
}

# Load template PPTX
template_path = "kad_penghargaan.pptx"
prs = Presentation(template_path)

# Pilih slide pertama
slide = prs.slides[0]

# Tambah text tengah (bawah Congratulation)
text = f"""Nama: {client_data['nama']}
IC: {mask_ic(client_data['ic'])}
No WhatsApp: {client_data['whatsapp']}
Email: {client_data['email']}
ID: {client_data['id']}
"""

left = Pt(340)
top = Pt(210)
width = Pt(300)
height = Pt(120)

txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = text
p.font.size = Pt(14)
p.font.bold = True
p.font.name = 'Arial'
p.font.color.rgb = prs.slides[0].shapes[1].text_frame.paragraphs[0].font.color.rgb  # sama warna

# Simpan sebagai fail baru
output_path = f"kad_{client_data['nama'].replace(' ', '_')}.pptx"
prs.save(output_path)
print("Kad penghargaan berjaya dijana:", output_path)
