# lock 2_Client.py v1.6.1 (cross-platform stable full)
import streamlit as st
import pandas as pd
import os, urllib.parse, json, time, shutil, subprocess, platform
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Conditional import untuk Windows sahaja
if platform.system() == "Windows":
    import comtypes.client
    import pythoncom
else:
    comtypes = None
    pythoncom = None

CSV_FILE = "pemilih.csv"
SETUP_FILE = "setup.json"

FOLDER_PPTX = "pptx"
FOLDER_SLIP = "slip"
os.makedirs(FOLDER_PPTX, exist_ok=True)
os.makedirs(FOLDER_SLIP, exist_ok=True)

# ====== STATE ======
if "busy" not in st.session_state:
    st.session_state["busy"] = False
if "uploads_locked" not in st.session_state:
    st.session_state["uploads_locked"] = False
if "ic_depan_bytes" not in st.session_state:
    st.session_state["ic_depan_bytes"] = None
if "ic_belakang_bytes" not in st.session_state:
    st.session_state["ic_belakang_bytes"] = None

# ====== HELPERS ======
def load_setup():
    with open(SETUP_FILE, "r") as f:
        return json.load(f)

def mask_ic(ic: str) -> str:
    if len(ic) == 12:
        return ic[:2] + "***" + ic[5:8] + "***" + ic[-1]
    return ic

def render_share_buttons(nama: str, masked_ic: str, profile_url: str):
    text = f"Kad Penghargaan BN\nNama: {nama}\nIC: {masked_ic}\nSemak: {profile_url}"
    enc_text = urllib.parse.quote(text)
    enc_url  = urllib.parse.quote(profile_url)
    wa = f"https://wa.me/?text={enc_text}"
    tg = f"https://t.me/share/url?url={enc_url}&text={enc_text}"
    fb = f"https://www.facebook.com/sharer/sharer.php?u={enc_url}"
    xx = f"https://twitter.com/intent/tweet?text={enc_text}&url={enc_url}"
    c1, c2, c3, c4 = st.columns(4)
    c1.link_button("🟢 WhatsApp", wa)
    c2.link_button("🔵 Telegram", tg)
    c3.link_button("🔷 Facebook", fb)
    c4.link_button("⚫ X (Twitter)", xx)

def _convert_pptx_to_pdf_headless(output_pptx: str, output_pdf: str) -> bool:
    candidates = [
        shutil.which("soffice"),
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    soffice = next((p for p in candidates if p and os.path.exists(p)), None)
    if not soffice:
        return False
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf",
             os.path.abspath(output_pptx),
             "--outdir", os.path.abspath(FOLDER_SLIP)],
            check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL
        )
        return os.path.exists(output_pdf)
    except Exception:
        return False

def _convert_pptx_to_pdf_com(output_pptx: str, output_pdf: str) -> bool:
    if not comtypes or not pythoncom:
        return False
    deck = None
    ppt = None
    try:
        pythoncom.CoInitialize()
        ppt = comtypes.client.CreateObject("PowerPoint.Application")
        ppt.Visible = True
        try:
            ppt.DisplayAlerts = 0
            ppt.WindowState = 2
        except:
            pass
        try:
            deck = ppt.Presentations.Open(os.path.abspath(output_pptx), WithWindow=False)
        except Exception:
            deck = ppt.Presentations.Open(os.path.abspath(output_pptx), WithWindow=True)
        ppSaveAsPDF = 32
        deck.SaveAs(os.path.abspath(output_pdf), ppSaveAsPDF)
        return os.path.exists(output_pdf)
    except Exception as e:
        print("❌ Gagal convert ke PDF (COM):", e)
        return False
    finally:
        try:
            if deck: deck.Close()
        except:
            pass
        try:
            if ppt: ppt.Quit()
        except:
            pass
        try:
            pythoncom.CoUninitialize()
        except:
            pass

def generate_kad_penghargaan(nama, ic_masked, whatsapp, email, bil):
    template_pptx = "kad_penghargaan.pptx"
    output_pptx = f"{FOLDER_PPTX}/output_kad_penghargaan_{nama.replace(' ', '_').lower()}.pptx"
    output_pdf = output_pptx.replace(".pptx", ".pdf").replace(FOLDER_PPTX, FOLDER_SLIP)

    if not os.path.exists(template_pptx):
        return None

    prs = Presentation(template_pptx)
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame and "your text here" in shape.text.strip().lower():
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"{nama.upper()}\nIC: {ic_masked}\nWhatsApp: {whatsapp}\nEmail: {email}\nNo Pendaftaran: {bil}"
            p.font.name = "Aptos Narrow"
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 32, 96)
            p.alignment = 1
            break
    prs.save(output_pptx)
    os.makedirs(FOLDER_SLIP, exist_ok=True)

    if _convert_pptx_to_pdf_headless(output_pptx, output_pdf):
        return output_pdf
    if platform.system() == "Windows" and comtypes:
        if _convert_pptx_to_pdf_com(output_pptx, output_pdf):
            return output_pdf
    return None

def _ui_lock_overlay():
    if st.session_state.get("busy"):
        st.markdown("""
        <style>
        [data-testid="stSidebar"], [data-testid="stHeader"], [data-testid="stToolbar"],
        [data-testid="stAppViewContainer"], [data-testid="stMain"] * {
            pointer-events: none !important;
        }
        #ui-lock {
            position: fixed; inset: 0; z-index: 99999;
            background: rgba(0,0,0,0.35);
            display: flex; align-items: center; justify-content: center;
            pointer-events: all !important;
        }
        .ui-card {
            background: white; border-radius: 14px; padding: 16px 22px;
            box-shadow: 0 10px 30px rgba(0,0,0,.25); font-size: 1.05rem;
        }
        </style>
        <div id="ui-lock"><div class="ui-card">⏳ Sedang memproses... Sila tunggu sampai selesai.</div></div>
        """, unsafe_allow_html=True)

# === UI ===
st.title("📝 Borang Pengundi (Client)")
__ = _ui_lock_overlay()

setup = load_setup()
dun = st.selectbox("DUN", setup.get("dun", []), disabled=st.session_state["busy"])
daerah = st.selectbox("Daerah Mengundi", setup.get("daerah_mengundi", []), disabled=st.session_state["busy"])
lokaliti = st.selectbox("Lokaliti", setup.get("lokaliti", []), disabled=st.session_state["busy"])

nama = st.text_input("Nama Penuh", disabled=st.session_state["busy"])
no_kp = st.text_input("No Kad Pengenalan (12 digit, tanpa -)", disabled=st.session_state["busy"])

upload_disabled = st.session_state["uploads_locked"] or st.session_state["busy"]
ic_depan = st.file_uploader("📷 Upload IC Depan", type=["jpg","jpeg","png"], disabled=upload_disabled, key="ic_depan_upl")
ic_belakang = st.file_uploader("📷 Upload IC Belakang", type=["jpg","jpeg","png"], disabled=upload_disabled, key="ic_belakang_upl")

if not st.session_state["uploads_locked"] and not st.session_state["busy"]:
    if ic_depan is not None:
        st.session_state["ic_depan_bytes"] = bytes(ic_depan.getbuffer())
    if ic_belakang is not None:
        st.session_state["ic_belakang_bytes"] = bytes(ic_belakang.getbuffer())
    if st.session_state["ic_depan_bytes"] and st.session_state["ic_belakang_bytes"]:
        st.session_state["uploads_locked"] = True

if st.session_state["uploads_locked"]:
    st.caption("🔒 Gambar IC Depan & Belakang dikunci (upload selesai).")

whatsapp = st.text_input("No WhatsApp", disabled=st.session_state["busy"])
email = st.text_input("Email", disabled=st.session_state["busy"])

# ====== Validation ======
name_ok = bool(nama.strip())
ic_input = no_kp.strip()
ic_ok = ic_input.isdigit() and len(ic_input) == 12
wa_ok = bool(whatsapp.strip())
both_uploaded_ok = (st.session_state["ic_depan_bytes"] and st.session_state["ic_belakang_bytes"])

ic_exists_in_csv = False
if os.path.exists(CSV_FILE) and ic_ok:
    _df = pd.read_csv(CSV_FILE, dtype={"no_kp": str})
    _kp_clean = _df["no_kp"].astype(str).str.replace("^'", "", regex=True)
    ic_exists_in_csv = ic_input in _kp_clean.values
if ic_exists_in_csv:
    st.warning("⚠️ No KP sudah wujud dalam sistem. Sila semak di halaman Profile.")

disable_submit = (
    st.session_state["busy"] or
    (not name_ok) or (not ic_ok) or (not both_uploaded_ok) or (not wa_ok) or ic_exists_in_csv
)

# ====== Form ======
with st.form("borang_pendaftaran", clear_on_submit=False):
    col1, col2 = st.columns(2)
    with col1:
        status = st.radio("Status", ["A (Aktif)", "P (Pasif)"], horizontal=True, disabled=st.session_state["busy"])
        sikap = st.radio("Sikap", ["P (Penyokong)", "G (Goyang)", "W (Wait & See)"], horizontal=True, disabled=st.session_state["busy"])
        umno = st.radio("UMNO", ["Y (Ya)", "N (Tidak)", "G (Gantung)"], horizontal=True, disabled=st.session_state["busy"])
        prbm = st.selectbox("PRBM", ["Ahli Biasa", "Veteran", "Ahli Seumur Hidup"], disabled=st.session_state["busy"])
    with col2:
        jawatan = st.text_input("Jawatan PDM", disabled=st.session_state["busy"])
        penilaian = st.text_input("Penilaian", disabled=st.session_state["busy"])
        blok = st.text_input("Blok", disabled=st.session_state["busy"])
        psywar = st.text_input("Psywar", disabled=st.session_state["busy"])
    submit = st.form_submit_button("💾 Simpan Data", disabled=disable_submit)

# ====== Submit Logic ======
if submit and not st.session_state["busy"]:
    st.session_state["busy"] = True
    _ui_lock_overlay()
    with st.spinner("⏳ Sila tunggu, sedang menyimpan & menjana kad..."):
        if not (ic_input.isdigit() and len(ic_input) == 12):
            st.error("❌ No KP mesti 12 digit")
            st.session_state["busy"] = False
            st.stop()
        if not nama.strip():
            st.error("❌ Nama tidak boleh kosong")
            st.session_state["busy"] = False
            st.stop()

        if os.path.exists(CSV_FILE):
            df = pd.read_csv(CSV_FILE, dtype={"no_kp": str})
        else:
            df = pd.DataFrame(columns=[
                "dun","daerah_mengundi","lokaliti","nama","no_kp","status","sikap","umno","prbm",
                "jawatan_pdm","penilaian","blok","psywar","ic_depan","ic_belakang","whatsapp","email"
            ])

        df["_kp_clean"] = df["no_kp"].astype(str).str.replace("^'", "", regex=True)
        if ic_input in df["_kp_clean"].values:
            st.warning("⚠️ No KP sudah wujud!")
            st.session_state["busy"] = False
            st.stop()

        os.makedirs("uploads", exist_ok=True)
        depan_path  = f"uploads/{ic_input}_front.jpg"   if st.session_state["ic_depan_bytes"] else ""
        belakang_path = f"uploads/{ic_input}_back.jpg"  if st.session_state["ic_belakang_bytes"] else ""
        if st.session_state["ic_depan_bytes"]:
            with open(depan_path, "wb") as f: f.write(st.session_state["ic_depan_bytes"])
        if st.session_state["ic_belakang_bytes"]:
            with open(belakang_path, "wb") as f: f.write(st.session_state["ic_belakang_bytes"])

        new_row = [
            dun, daerah, lokaliti, nama,
            f"'{ic_input}", status, sikap, umno, prbm,
            jawatan, penilaian, blok, psywar,
            depan_path, belakang_path,
            f"'{whatsapp}", email
        ]

        expected_cols = [
            "dun","daerah_mengundi","lokaliti","nama","no_kp","status","sikap","umno","prbm",
            "jawatan_pdm","penilaian","blok","psywar","ic_depan","ic_belakang","whatsapp","email"
        ]
        df = df.reindex(columns=expected_cols)
        df.loc[len(df)] = new_row
        df.to_csv(CSV_FILE, index=False)

        masked_ic = mask_ic(ic_input)
        if os.path.exists("Logo-BN.png"):
            st.image("Logo-BN.png", width=120)
        st.success(f"✅ Data berjaya disimpan!\n\nNama: {nama}\nIC: {masked_ic}\nWhatsApp: {whatsapp}\nEmail: {email}")

        kad_file = generate_kad_penghargaan(nama, masked_ic, whatsapp, email, bil=len(df))
        if kad_file and os.path.exists(kad_file):
            with open(kad_file, "rb") as f:
                st.download_button("⬇️ Muat Turun Kad PDF", f, file_name=os.path.basename(kad_file), mime="application/pdf")

            public_base = setup.get("public_base_url", "http://localhost:8501")
            profile_url = f"{public_base}/Profile?ic={ic_input}"
            st.info("📣 Kongsi kad ini:")
            render_share_buttons(nama, masked_ic, profile_url)
        else:
            st.error("❌ Gagal jana Kad Penghargaan (PDF). Pastikan PowerPoint/LibreOffice terpasang dengan betul.")

        if whatsapp:
            msg = urllib.parse.quote(
                f"✅ Pendaftaran berjaya!\nNama: {nama}\nIC: {masked_ic}\nWhatsApp: {whatsapp}\nEmail: {email}"
            )
            wa_link = f"https://wa.me/{whatsapp}?text={msg}"
            st.markdown(f"[📲 Hantar Slip ke WhatsApp]({wa_link})", unsafe_allow_html=True)

    st.session_state["uploads_locked"] = False
    st.session_state["ic_depan_bytes"] = None
    st.session_state["ic_belakang_bytes"] = None
    st.session_state["busy"] = False
    __ = _ui_lock_overlay()
